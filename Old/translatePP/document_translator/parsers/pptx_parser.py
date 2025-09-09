from typing import Dict, Any
from pptx import Presentation
import json

class PPTXParser:
    def extract_content(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text and formatting information from a PowerPoint file.
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            Dictionary containing text and formatting information
        """
        prs = Presentation(file_path)
        text_content = []
        formatting_info = []
        
        for slide_number, slide in enumerate(prs.slides):
            slide_content = {
                "slide_number": slide_number + 1,
                "shapes": []
            }
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape_content = {
                        "text": shape.text,
                        "formatting": {
                            "font_name": shape.text_frame.paragraphs[0].font.name if shape.text_frame.paragraphs else None,
                            "font_size": shape.text_frame.paragraphs[0].font.size if shape.text_frame.paragraphs else None,
                            "is_bold": shape.text_frame.paragraphs[0].font.bold if shape.text_frame.paragraphs else None,
                            "is_italic": shape.text_frame.paragraphs[0].font.italic if shape.text_frame.paragraphs else None,
                            "alignment": shape.text_frame.paragraphs[0].alignment if shape.text_frame.paragraphs else None
                        }
                    }
                    text_content.append(shape.text)
                    slide_content["shapes"].append(shape_content)
            
            formatting_info.append(slide_content)
        
        return {
            "text": "\n".join(text_content),
            "formatting": formatting_info
        }
    
    def create_document(self, file_path: str, translated_text: str, 
                       formatting_info: Dict[str, Any]) -> None:
        """
        Create a new PowerPoint file with translated content and preserved formatting.
        
        Args:
            file_path: Path where the new PowerPoint file should be saved
            translated_text: The translated text content
            formatting_info: Dictionary containing formatting information
        """
        prs = Presentation()
        
        # Split translated text into chunks (one per shape)
        translated_chunks = translated_text.split("\n")
        chunk_index = 0
        
        for slide_info in formatting_info:
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
            
            for shape_info in slide_info["shapes"]:
                if chunk_index < len(translated_chunks):
                    shape = slide.shapes.add_textbox(
                        shape_info.get("left", 0),
                        shape_info.get("top", 0),
                        shape_info.get("width", 100),
                        shape_info.get("height", 100)
                    )
                    
                    text_frame = shape.text_frame
                    paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
                    
                    # Apply text
                    paragraph.text = translated_chunks[chunk_index]
                    
                    # Apply formatting
                    if shape_info["formatting"]:
                        font = paragraph.font
                        if shape_info["formatting"]["font_name"]:
                            font.name = shape_info["formatting"]["font_name"]
                        if shape_info["formatting"]["font_size"]:
                            font.size = shape_info["formatting"]["font_size"]
                        if shape_info["formatting"]["is_bold"] is not None:
                            font.bold = shape_info["formatting"]["is_bold"]
                        if shape_info["formatting"]["is_italic"] is not None:
                            font.italic = shape_info["formatting"]["is_italic"]
                        if shape_info["formatting"]["alignment"]:
                            paragraph.alignment = shape_info["formatting"]["alignment"]
                    
                    chunk_index += 1
        
        prs.save(file_path) 